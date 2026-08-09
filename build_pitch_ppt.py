#!/usr/bin/env python3.11
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DEEP  = RGBColor(0x0F,0x39,0x54)
DARK  = RGBColor(0x16,0x21,0x3E)
GOLD  = RGBColor(0xC9,0xA8,0x4C)
LGOLD = RGBColor(0xF4,0xD0,0x6F)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LIGHT = RGBColor(0xE2,0xE8,0xF0)
GREY  = RGBColor(0x94,0xA3,0xB8)
CARD  = RGBColor(0x1A,0x28,0x42)
CARD2 = RGBColor(0x1E,0x2C,0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def box(slide, l,t,w,h, fill=None, line=None, line_w=None, radius=True, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l,t,w,h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    if not shadow: shp.shadow.inherit = False
    return shp

def txt(slide, l,t,w,h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l,t,w,h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=Pt(5); tf.margin_right=Pt(5); tf.margin_top=Pt(3); tf.margin_bottom=Pt(3)
    for i,(s,size,color,bold,*rest) in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(rest[0] if rest else 3)
        if len(rest)>1 and rest[1]: p.space_before = Pt(rest[1])
        r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Microsoft YaHei"
    return tb

def title(slide, text, sub=None):
    box(slide, 0, 0, SW, Inches(1.1), fill=DARK)
    box(slide, 0, Inches(1.1), SW, Pt(3), fill=GOLD)
    txt(slide, Inches(0.55), Inches(0.15), SW-Inches(1), Inches(0.85),
        [(text, 26, LGOLD, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        txt(slide, Inches(0.57), Inches(0.74), SW-Inches(1), Inches(0.35),
            [(sub, 12.5, GREY, False)])

def section_num(slide, n):
    box(slide, Inches(0.55), Inches(0.18), Inches(0.55), Inches(0.55), fill=GOLD, radius=True)
    txt(slide, Inches(0.55), Inches(0.18), Inches(0.55), Inches(0.55), [(n,20,DEEP,True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 1: Cover =====
s = prs.slides.add_slide(BLANK); bg(s, DEEP)
box(s, Inches(0.9), Inches(2.4), Inches(0.14), Inches(1.7), fill=GOLD)
box(s, 0, Inches(2.4), SW, Pt(2), fill=GOLD)
txt(s, Inches(1.2), Inches(2.5), Inches(11), Inches(1.0), [("一四油 · 招商推介", 42, WHITE, True)])
txt(s, Inches(1.2), Inches(3.55), Inches(11), Inches(0.7), [("大健康轻创业 · 三级进阶 · 多劳多得", 20, LGOLD, False)])
txt(s, Inches(1.2), Inches(5.6), Inches(11), Inches(0.8),
    [("低门槛起步 · 完整赋能体系 · 透明收益分配", 15, LIGHT, False),
     ("内部培训资料 · 数据以官方最新发布为准", 12, GREY, False, 4)])

# ===== Slide 2: Hook =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "为什么是现在的机会", "一个低门槛、可复制的大健康创业模型")
stats = [("19440 元","起步门槛（3 箱 144 瓶）",GOLD),
         ("135 元/瓶","经销商永久批发价",GOLD),
         ("25–85 元","真实单瓶利润空间",GOLD),
         ("3 级","清晰可量化的进阶通道",GOLD)]
x = Inches(0.7); w = Inches(2.85); gap = Inches(0.25)
for i,(big,small,c) in enumerate(stats):
    bx = x + i*(w+gap)
    box(s, bx, Inches(1.55), w, Inches(1.7), fill=CARD, line=GOLD, line_w=Pt(1))
    txt(s, bx, Inches(1.7), w, Inches(0.95), [(big, 30, c, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, bx, Inches(2.65), w, Inches(0.55), [(small, 12.5, LIGHT, False)], align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.5), [("核心逻辑：用零售端的价盘空间，构建可裂变的经销网络", 16, LGOLD, True)])
box(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(2.2), fill=CARD, line=GOLD, line_w=Pt(1), radius=True)
txt(s, Inches(0.95), Inches(4.5), Inches(11.4), Inches(2.0),
    [("● 产品端：复购型健康消费，用户黏性强、复购周期稳定", 14, LIGHT, False, 7),
     ("● 渠道端：3 箱即可成为经销商，门槛低、决策快", 14, LIGHT, False, 7),
     ("● 收益端：单瓶利润空间 25–191 元，随等级跃升而放大", 14, LIGHT, False, 7),
     ("● 团队端：下级进货持续产生服务费与补贴，越带越宽", 14, LIGHT, False, 7)])

# ===== Slide 3: 产品简介 (占位) =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "产品简介 · 一四油", "【待补充：以下为结构框架，请提供产品资料后填实】")
box(s, Inches(0.7), Inches(1.5), Inches(5.85), Inches(4.9), fill=CARD, line=GOLD, line_w=Pt(1))
txt(s, Inches(0.95), Inches(1.65), Inches(5.4), Inches(0.5), [("◆ 产品定位与核心成分", 15, LGOLD, True)])
txt(s, Inches(0.95), Inches(2.2), Inches(5.4), Inches(4.0),
    [("• 产品名称 / 形态：________", 13, LIGHT, False, 6),
     ("• 核心成分：________", 13, LIGHT, False, 6),
     ("• 核心功效（合规表述）：________", 13, LIGHT, False, 6),
     ("• 适用人群：________", 13, LIGHT, False, 6),
     ("• 与同类产品的差异化：________", 13, LIGHT, False, 6)])
box(s, Inches(6.75), Inches(1.5), Inches(5.85), Inches(4.9), fill=CARD, line=GOLD, line_w=Pt(1))
txt(s, Inches(7.0), Inches(1.65), Inches(5.4), Inches(0.5), [("◆ 用户价值 & 市场痛点", 15, LGOLD, True)])
txt(s, Inches(7.0), Inches(2.2), Inches(5.4), Inches(4.0),
    [("• 目标用户痛点：________", 13, LIGHT, False, 6),
     ("• 使用场景：________", 13, LIGHT, False, 6),
     ("• 复购周期 / 客单价：________", 13, LIGHT, False, 6),
     ("• 主推话术 / 体验装设计：________", 13, LIGHT, False, 6),
     ("• 资质 / 检测报告 / 背书：________", 13, LIGHT, False, 6)])

# ===== Slide 4: 价盘与利润空间 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "价盘与利润空间 · 一目了然", "采购数量不累计 · 后续进货永久享受对应等级单价")
rows = [("采购规格","总价","单瓶","身份/适用"),
        ("1 瓶","286 元","286 元","零售"),
        ("3 瓶","660 元","220 元","小批量"),
        ("24 瓶（半箱）","4560 元","190 元","大批量"),
        ("48 瓶（1 箱）","7680 元","160 元","会员价档"),
        ("3 箱（144 瓶）","19440 元","135 元","★经销商"),
        ("18 箱（864 瓶）","99360 元","115 元","★大经销商"),
        ("95 箱（4560 瓶）","433200 元","95 元","★总经销商")]
gt = s.shapes.add_table(len(rows),4, Inches(0.7), Inches(1.5), Inches(8.2), Inches(4.6)).table
gt.columns[0].width=Inches(2.5); gt.columns[1].width=Inches(1.7); gt.columns[2].width=Inches(1.4); gt.columns[3].width=Inches(2.6)
for ci in range(4):
    c=gt.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb=DEEP
    p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=rows[0][ci]; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=LGOLD
for ri in range(1,len(rows)):
    dealer = ri>=5
    for ci in range(4):
        c=gt.cell(ri,ci); c.fill.solid(); c.fill.fore_color.rgb = CARD2 if dealer else CARD
        p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if ci<3 else PP_ALIGN.LEFT
        r=p.add_run(); r.text=rows[ri][ci]
        r.font.size=Pt(12); r.font.bold=(ci==2 or dealer); r.font.color.rgb=LGOLD if (ci==2 or dealer) else LIGHT
        c.margin_top=Pt(2); c.margin_bottom=Pt(2)
# right panel: 利润空间
box(s, Inches(9.2), Inches(1.5), Inches(3.5), Inches(4.6), fill=CARD, line=GOLD, line_w=Pt(1))
txt(s, Inches(9.4), Inches(1.65), Inches(3.1), Inches(0.5), [("单瓶利润空间", 15, LGOLD, True)])
txt(s, Inches(9.4), Inches(2.25), Inches(3.1), Inches(3.7),
    [("零售 286 → 经销 135", 12.5, LIGHT, False, 5),
     ("  单瓶空间 151 元", 12.5, LGOLD, True, 10),
     ("会员价 160 → 经销 135", 12.5, LIGHT, False, 5),
     ("  单瓶赚 25 元", 12.5, LGOLD, True, 10),
     ("24瓶档 190 → 经销 135", 12.5, LIGHT, False, 5),
     ("  单瓶赚 55 元", 12.5, LGOLD, True, 10),
     ("3瓶档 220 → 经销 135", 12.5, LIGHT, False, 5),
     ("  单瓶赚 85 元", 12.5, LGOLD, True, 10)])

# ===== Slide 5: 传销辨析 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "正规经销 vs 传销 · 一图看清", "我们卖产品赚差价，不是拉人头分钱")
# two columns
box(s, Inches(0.7), Inches(1.5), Inches(5.85), Inches(4.7), fill=CARD, line=RGBColor(0x8B,0x3A,0x3A), line_w=Pt(1.5))
txt(s, Inches(0.95), Inches(1.65), Inches(5.4), Inches(0.5), [("✗ 传销的特征（我们不是）", 15, RGBColor(0xE8,0x8A,0x8A), True)])
txt(s, Inches(0.95), Inches(2.25), Inches(5.4), Inches(3.8),
    [("• 靠拉人头、交入门费赚钱", 13.5, LIGHT, False, 9),
     ("• 没有真实产品 / 产品只是道具", 13.5, LIGHT, False, 9),
     ("• 承诺躺赚、静态返利", 13.5, LIGHT, False, 9),
     ("• 多级抽成、无退货保障", 13.5, LIGHT, False, 9),
     ("• 收入来自下线人头，不靠卖货", 13.5, LIGHT, False, 9)])
box(s, Inches(6.75), Inches(1.5), Inches(5.85), Inches(4.7), fill=CARD, line=GOLD, line_w=Pt(1.5))
txt(s, Inches(7.0), Inches(1.65), Inches(5.4), Inches(0.5), [("✓ 我们的模式（正规经销）", 15, LGOLD, True)])
txt(s, Inches(7.0), Inches(2.25), Inches(5.4), Inches(3.8),
    [("• 进货真实产品，靠卖差价获利", 13.5, LIGHT, False, 9),
     ("• 产品可自用 / 复购 / 真实消费", 13.5, LIGHT, False, 9),
     ("• 首次进货 1 月内可无条件退", 13.5, LIGHT, False, 9),
     ("• 收益来自销售与服务，多劳多得", 13.5, LIGHT, False, 9),
     ("• 需注册公司、对公纳税（总经销）", 13.5, LIGHT, False, 9)])
box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7), fill=DEEP, line=GOLD, line_w=Pt(1))
txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
    [("一句话：你赚的每一分钱，都来自真实的产品销售与售后服务，不是下线的人头费。", 13.5, LGOLD, True)],
    anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 6: 三级进阶路径 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "三级进阶路径 · 清晰可量化", "每升一级，批发价下降、单瓶利润放大")
tiers=[("① 基础经销商","一次性 3 箱 144 瓶","19440 元","135 元/瓶"),
       ("② 大经销商","当月团队累计 18 箱","864 瓶","115 元/瓶"),
       ("③ 总经销商","当月团队累计 95 箱","4560 瓶","95 元/瓶")]
x=Inches(0.7); w=Inches(3.8); gap=Inches(0.45)
for i,(name,cond,vol,price) in enumerate(tiers):
    bx=x+i*(w+gap)
    box(s, bx, Inches(1.7), w, Inches(3.2), fill=CARD, line=GOLD, line_w=Pt(1.5))
    txt(s, bx, Inches(1.9), w, Inches(0.6), [(name,17,LGOLD,True)], align=PP_ALIGN.CENTER)
    txt(s, bx, Inches(2.65), w, Inches(0.9),
        [("晋级："+cond, 13, LIGHT, False, 4),("（"+vol+"）",13,LIGHT,False)], align=PP_ALIGN.CENTER)
    txt(s, bx, Inches(3.75), w, Inches(0.8), [(price, 26, GOLD, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i<2:
        txt(s, bx+w-Inches(0.05), Inches(3.0), Inches(0.5), Inches(0.6), [("→",24,GOLD,True)], align=PP_ALIGN.CENTER)
box(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.1), fill=CARD2, line=GOLD, line_w=Pt(1))
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.9),
    [("关键规则：采购数量不累计 · 后续进货永久享受晋级后单价 · 月度计算周期为当月 26 日—次月 25 日", 13.5, LIGHT, False)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 6: 收益模型 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "收益模型 · 每一级能赚多少", "收益 = 差价 + 服务费 + 进货补贴 + 年度分红（多劳多得）")
cols=[("基础经销商","135 元/瓶 进货",
       ["售 144 瓶（按零售价）可赚 12240 元","对应毛利率约 63%","实际按会员/小批量档单瓶赚 25–85 元","服务费 8 元/瓶（下级进货）","进货补贴 3%–23%"],GOLD),
      ("大经销商","115 元/瓶 进货",
       ["对比会员价单瓶赚 45 元","对比零售价单瓶赚 171 元","服务费 8 元/瓶","进货补贴 3%–23%"],GOLD),
      ("总经销商","95 元/瓶 进货",
       ["对比会员价单瓶赚 65 元","对比零售价单瓶赚 191 元","服务费 10 元/瓶（公司直付）","进货补贴 3%–23% 上不封顶"],GOLD)]
x=Inches(0.7); w=Inches(3.8); gap=Inches(0.45)
for i,(name,price,items,c) in enumerate(cols):
    bx=x+i*(w+gap)
    box(s, bx, Inches(1.55), w, Inches(4.6), fill=CARD, line=GOLD, line_w=Pt(1))
    txt(s, bx, Inches(1.7), w, Inches(0.5), [(name,16,LGOLD,True)], align=PP_ALIGN.CENTER)
    txt(s, bx, Inches(2.25), w, Inches(0.5), [(price,15,c,True)], align=PP_ALIGN.CENTER)
    txt(s, bx+Inches(0.3), Inches(2.95), w-Inches(0.5), Inches(3.0),
        [("• "+t, 13, LIGHT, False, 9) for t in items])
box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7), fill=DEEP, line=GOLD, line_w=Pt(1))
txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
    [("案例：总经销团队月销售额 2000 万（约 20 万瓶）→ 服务费 200 万 + 进货补贴 120 万 + 年度分红 80 万 ≈ 400 万元/月", 13, LGOLD, True)],
    anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 8: 新人 90 天真实路径 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "新人前 90 天 · 真实路径", "不画大饼——回本 + 微利，是第一步目标")
months=[("第 1 月","自用好 + 送 10 位精准亲友体验","转化 5 个复购客户，练熟产品话术",GOLD),
        ("第 2 月","老客转介绍 + 社群真实分享","累计 24 个精准客户，完成铺货动作",GOLD),
        ("第 3 月","带出 1–2 个同频伙伴","开始有服务费收入，团队起量",GOLD)]
x=Inches(0.7); w=Inches(3.85); gap=Inches(0.3)
for i,(m,a,b,c) in enumerate(months):
    bx=x+i*(w+gap)
    box(s, bx, Inches(1.6), w, Inches(3.4), fill=CARD, line=GOLD, line_w=Pt(1.5))
    txt(s, bx, Inches(1.8), w, Inches(0.6), [(m,18,c,True)], align=PP_ALIGN.CENTER)
    txt(s, bx+Inches(0.25), Inches(2.6), w-Inches(0.5), Inches(2.2),
        [("• "+a, 13, LIGHT, False, 8),("• "+b, 13, LIGHT, False, 8)])
box(s, Inches(0.7), Inches(5.3), Inches(11.9), Inches(1.1), fill=CARD2, line=GOLD, line_w=Pt(1))
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(0.9),
    [("现实预期：前 3 个月以「回本 + 微利」为主，团队与被动收入从第 3 月起逐步起量。", 14, LGOLD, True, 4),
     ("稳扎稳打，好过一夜暴富——400 万/月的案例是塔尖，不是起点。", 12.5, LIGHT, False)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 9: 赋能体系 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "赋能体系 · 你不是一个人", "从培训到售后，标准化流程兜底")
items=[("系统培训","产品知识 · 铺货技巧 · 回款话术 · 调理反应应对"),
       ("标准售后","7 天体验打卡 + 1 个月跟踪 + 3 个月回访"),
       ("团队扶持","大经销筛选 6 名核心伙伴，复制起步模式"),
       ("透明分配","差价/服务费/补贴按规则公开透明发放"),
       ("退货保障","首次进货 1 月内未使用可无条件退"),
       ("合规经营","总经销需注册公司、对公转账、合法纳税")]
x0=Inches(0.7); y0=Inches(1.55); w=Inches(3.85); h=Inches(1.55); gx=Inches(0.3); gy=Inches(0.25)
for i,(t,d) in enumerate(items):
    r=i//3; c=i%3
    bx=x0+c*(w+gx); by=y0+r*(h+gy)
    box(s, bx, by, w, h, fill=CARD, line=GOLD, line_w=Pt(1))
    txt(s, bx+Inches(0.2), by+Inches(0.15), w-Inches(0.4), Inches(0.5), [(t,15,LGOLD,True)])
    txt(s, bx+Inches(0.2), by+Inches(0.7), w-Inches(0.4), Inches(0.8), [(d,12.5,LIGHT,False)])

# ===== Slide 8: 适合谁 / 起步 =====
s = prs.slides.add_slide(BLANK); bg(s, DARK); title(s, "谁适合做 · 起步三步", "无论身份，都能找到自己的切入点")
box(s, Inches(0.7), Inches(1.5), Inches(5.7), Inches(4.9), fill=CARD, line=GOLD, line_w=Pt(1))
txt(s, Inches(0.95), Inches(1.65), Inches(5.2), Inches(0.5), [("◆ 适合人群", 15, LGOLD, True)])
txt(s, Inches(0.95), Inches(2.2), Inches(5.2), Inches(4.0),
    [("• 健康行业从业者 / 养生馆主", 13.5, LIGHT, False, 7),
     ("• 宝妈 / 退休人群（时间灵活）", 13.5, LIGHT, False, 7),
     ("• 微商 / 社群团长（有私域流量）", 13.5, LIGHT, False, 7),
     ("• 企业主 / 销售背景（带团队强）", 13.5, LIGHT, False, 7),
     ("• 认同产品、想轻创业的普通人", 13.5, LIGHT, False, 7)])
box(s, Inches(6.6), Inches(1.5), Inches(6.0), Inches(4.9), fill=CARD, line=GOLD, line_w=Pt(1))
txt(s, Inches(6.85), Inches(1.65), Inches(5.5), Inches(0.5), [("◆ 起步三步", 15, LGOLD, True)])
txt(s, Inches(6.85), Inches(2.2), Inches(5.5), Inches(4.0),
    [("① 采购 3 箱（144 瓶）19440 元 → 成为经销商", 13.5, LIGHT, False, 9),
     ("② 参加系统培训 + 完成 24 个精准铺货", 13.5, LIGHT, False, 9),
     ("③ 带团队冲 18 箱→95 箱，晋级大/总经销", 13.5, LIGHT, False, 9),
     ("全程：差价 + 服务费 + 补贴 + 分红，多劳多得", 12.5, LGOLD, True, 9)])

# ===== Slide 9: CTA =====
s = prs.slides.add_slide(BLANK); bg(s, DEEP)
box(s, 0, Inches(1.4), SW, Pt(3), fill=GOLD)
txt(s, Inches(1), Inches(1.7), Inches(11.3), Inches(1.0), [("现在，就是最好的起步时机", 34, WHITE, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(2.9), Inches(11.3), Inches(0.8),
    [("第一步只需：采购 3 箱（144 瓶）= 19440 元，即刻成为经销商", 18, LGOLD, False)], align=PP_ALIGN.CENTER)
box(s, Inches(2.4), Inches(4.0), Inches(8.5), Inches(1.5), fill=CARD, line=GOLD, line_w=Pt(1.5))
txt(s, Inches(2.6), Inches(4.15), Inches(8.1), Inches(1.2),
    [("联系我们，获取：产品资料 · 培训入口 · 1对1起步辅导", 15, LIGHT, False),
     ("（二维码 / 微信号 见下方）", 12, GREY, False, 5)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(5.4), Inches(5.8), Inches(2.5), Inches(1.0), fill=None, line=GOLD, line_w=Pt(1.5))
txt(s, Inches(5.4), Inches(5.8), Inches(2.5), Inches(1.0), [("二维码位", 13, GREY, False)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 10: 封底 =====
s = prs.slides.add_slide(BLANK); bg(s, DEEP)
box(s, 0, Inches(3.3), SW, Pt(3), fill=GOLD)
txt(s, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2), [("多劳多得 · 踩准节点 · 跟上节奏", 32, WHITE, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.8), Inches(11.3), Inches(0.6), [("一四油经销商制度参考 · 内部培训资料", 14, LGOLD, False)], align=PP_ALIGN.CENTER)

prs.save("/workspace/一四油招商推介.pptx")
print("Saved: /workspace/一四油招商推介.pptx, slides =", len(prs.slides._sldIdLst))
