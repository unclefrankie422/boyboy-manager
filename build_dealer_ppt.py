#!/usr/bin/env python3.11
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 颜色
DEEP   = RGBColor(0x0F,0x39,0x54)
DARK   = RGBColor(0x16,0x21,0x3E)
GOLD   = RGBColor(0xC9,0xA8,0x4C)
LGOLD  = RGBColor(0xF4,0xD0,0x6F)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
LIGHT  = RGBColor(0xE2,0xE8,0xF0)
GREY   = RGBColor(0x94,0xA3,0xB8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def box(slide, l,t,w,h, fill=None, line=None, line_w=None, radius=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l,t,w,h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp

def txt(slide, l,t,w,h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l,t,w,h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=Pt(4); tf.margin_right=Pt(4); tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    for i,(s,size,color,bold,*rest) in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        if rest and rest[0]: p.space_after = Pt(rest[0])
        else: p.space_after = Pt(2)
        if len(rest)>1 and rest[1]: p.space_before = Pt(rest[1])
        r = p.add_run(); r.text = s
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = "Microsoft YaHei"
    return tb

def title_bar(slide, text):
    box(slide, 0, 0, SW, Inches(1.05), fill=DARK)
    box(slide, 0, Inches(1.05), SW, Pt(3), fill=GOLD)
    txt(slide, Inches(0.5), Inches(0.12), SW-Inches(1), Inches(0.85),
        [(text, 26, LGOLD, True)], anchor=MSO_ANCHOR.MIDDLE)

# ---------- Slide 1: Cover ----------
s = prs.slides.add_slide(BLANK); bg(s, DEEP)
box(slide:=s, 0, Inches(2.55), SW, Pt(3), fill=GOLD)
txt(s, Inches(0.8), Inches(2.7), SW-Inches(1.6), Inches(1.4),
    [("一四油 · 经销商进阶计划", 40, WHITE, True)])
txt(s, Inches(0.8), Inches(3.9), SW-Inches(1.6), Inches(0.8),
    [("全档位价盘 · 三级制度 · 收益详解", 20, LGOLD, False)])
txt(s, Inches(0.8), Inches(5.7), SW-Inches(1.6), Inches(0.6),
    [("内部培训资料 · 数据以官方最新发布为准", 13, GREY, False)])
box(s, Inches(0.8), Inches(2.55), Inches(0.12), Inches(1.4), fill=GOLD)

# ---------- Slide 2: 价格明细 ----------
s = prs.slides.add_slide(BLANK); bg(s, DARK); title_bar(s, "一、全档位价格明细")
rows = [
    ("采购规格","总价","折合单瓶","适用人群 / 身份"),
    ("1 瓶","286 元","286 元","普通消费者 · 零售"),
    ("3 瓶","660 元","220 元","普通消费者 · 小批量"),
    ("24 瓶（半箱）","4560 元","190 元","普通消费者 · 大批量"),
    ("48 瓶（1 箱）","7680 元","160 元","会员价档位（普通消费者/入门经销商）"),
    ("3 箱（144 瓶）","19440 元","135 元","★ 经销商进货价 · 一次性采购即成为经销商"),
    ("18 箱（864 瓶）","99360 元","115 元","★ 大经销商进货价 · 当月团队累计进货晋升"),
    ("95 箱（4560 瓶）","433200 元","95 元","★ 总经销商进货价 · 当月团队累计进货晋升"),
]
nrows, ncols = len(rows), 4
gt = s.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.4)).table
gt.columns[0].width=Inches(2.7); gt.columns[1].width=Inches(2.0)
gt.columns[2].width=Inches(2.0); gt.columns[3].width=Inches(5.63)
for ci in range(ncols):
    c = gt.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb = DEEP
    p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=rows[0][ci]; r.font.size=Pt(14); r.font.bold=True; r.font.color.rgb=LGOLD
for ri in range(1,nrows):
    is_dealer = ri>=5
    for ci in range(ncols):
        c=gt.cell(ri,ci); c.fill.solid()
        c.fill.fore_color.rgb = RGBColor(0x1E,0x2C,0x4A) if is_dealer else RGBColor(0x1A,0x24,0x3A)
        p=c.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.CENTER if ci<3 else PP_ALIGN.LEFT
        r=p.add_run(); r.text=rows[ri][ci]
        r.font.size=Pt(12.5); r.font.bold = (ci==2 or is_dealer)
        r.font.color.rgb = LGOLD if (ci==2 or is_dealer) else LIGHT
        c.margin_top=Pt(3); c.margin_bottom=Pt(3)
txt(s, Inches(0.5), Inches(6.95), Inches(12), Inches(0.4),
    [("采购数量不累计 · 后续进货永久享受对应等级单价 · 月度周期：当月26日—次月25日", 11, GREY, False)])

# ---------- helper for tier slides ----------
def tier_slide(num, name, meta, rights, incomes, extra=None):
    s = prs.slides.add_slide(BLANK); bg(s, DARK); 
    title_bar(s, f"{num} {name}")
    # meta banner
    txt(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.55),
        [(meta, 15, LGOLD, True)])
    # two columns: 权利义务 / 收益权益
    colw = Inches(6.0)
    box(s, Inches(0.5), Inches(1.95), colw, Inches(4.9), fill=RGBColor(0x1A,0x28,0x42), line=GOLD, line_w=Pt(1), radius=True)
    box(s, Inches(6.85), Inches(1.95), colw, Inches(4.9), fill=RGBColor(0x1A,0x28,0x42), line=GOLD, line_w=Pt(1), radius=True)
    txt(s, Inches(0.7), Inches(2.05), colw-Inches(0.4), Inches(0.4), [("◆ 权利义务", 15, LGOLD, True)])
    txt(s, Inches(7.05), Inches(2.05), colw-Inches(0.4), Inches(0.4), [("◆ 收益权益", 15, LGOLD, True)])
    # build bullet text
    rl = [(("- "+t, 12.5, LIGHT, False, 4)) for t in rights]
    il = [(("- "+t, 12.5, LIGHT, False, 4)) for t in incomes]
    txt(s, Inches(0.7), Inches(2.5), colw-Inches(0.4), Inches(4.2), rl)
    txt(s, Inches(7.05), Inches(2.5), colw-Inches(0.4), Inches(4.2), il)
    if extra:
        box(s, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.45), fill=DEEP, radius=True)
        txt(s, Inches(0.7), Inches(6.98), Inches(12), Inches(0.4), [(extra, 11.5, LGOLD, True)], anchor=MSO_ANCHOR.MIDDLE)

# ---------- Slide 3: 基础经销商 ----------
tier_slide("①","基础经销商",
  "晋级：一次性购买 3 箱（144 瓶）共 19440 元 → 永久 135 元/瓶批发价",
  ["后续进货永久享受 135 元/瓶 批发价",
   "需参与系统培训：产品知识、铺货技巧、回款话术、调理反应应对",
   "需完成 24 个精准市场铺货，执行 7 天体验打卡 + 1 个月跟踪 + 3 个月回访标准售后",
   "首次进货 1 个月内未使用/未损坏可无条件退货；复购后仅质量问题可退"],
  ["差价：会员价 160 元售，单瓶赚 25 元；24 瓶档 190 元售，单瓶赚 55 元；3 瓶档 220 元售，单瓶赚 85 元",
   "售 144 瓶可赚 12240 元，毛利率约 63%",
   "服务费：分享同级经销商进货，获上级支付 8 元/瓶",
   "进货补贴：当月累计进货额 1 万–800 万，享 3%–23% 梯度补贴，上不封顶",
   "年度分红：团队年销售额 600 万–6000 万，享 3%–6%（需晋升总经销注册公司纳税后）",
   "特别贡献奖：按个人自购年销售额 10% 发，总额限定 100 万，发完即止"])

# ---------- Slide 4: 大经销商 ----------
tier_slide("②","大经销商",
  "晋级：当月团队累计进货 18 箱（864 瓶）晋升 → 永久 115 元/瓶批发价",
  ["后续进货永久享受 115 元/瓶 批发价",
   "需筛选 6 名认可产品、有创业意愿的核心伙伴，复制起步模式扶持下属市场",
   "需掌握团队搭建、裂变模式复制、团队管理、收益拆解技能",
   "需按规则给下级发放差价、服务费、进货补贴，公开透明分配"],
  ["差价：批发价 115 元，对比会员价最高单瓶赚 45 元，对比零售价单瓶赚 171 元",
   "服务费：分享同级经销商进货，获上级支付 8 元/瓶",
   "进货补贴：对应进货额 3%–23% 梯度补贴",
   "年度分红：团队年销售额达标可参与 3%–6%（需注册公司纳税后）",
   "特别贡献奖：按个人自购年销售额 10% 参与发放"])

# ---------- Slide 5: 总经销商 ----------
tier_slide("③","总经销商",
  "晋级：当月团队累计进货 95 箱（4560 瓶）晋升 → 永久 95 元/瓶批发价（需注册公司·对公转账·合法纳税·直接对接工厂）",
  ["后续进货永久享受 95 元/瓶 批发价",
   "注册公司经营范围：远程健康管理/健康咨询（不含诊疗）/养生保健（非医疗）/中医养生（非医疗）/日用百货/体育健康/预包装食品销售",
   "统筹团队全域运营，定向培育大经销商，推动团队规模裂变，统筹团队售后",
   "服务费由公司直接支付，需按规则给下级发放，不得违规截留",
   "不得在网上/线下乱价销售，否则不予退货且取消经销商资格"],
  ["差价：批发价 95 元，对比会员价最高单瓶赚 65 元，对比零售价单瓶赚 191 元",
   "服务费：由公司直接支付 10 元/瓶",
   "进货补贴：对应进货额 3%–23% 梯度补贴，上不封顶",
   "年度分红：团队年销售额 600 万–6000 万，享 3%–6%",
   "特别贡献奖：按个人自购年销售额 10% 发，总额限定 100 万，发完即止"],
  extra="收益案例：团队月销售额 2000 万（约 20 万瓶）→ 服务费 200 万 + 进货补贴 120 万 + 年度分红 80 万 = 三项合计约 400 万元/月")

# ---------- Slide 6: 通用规则 ----------
s = prs.slides.add_slide(BLANK); bg(s, DARK); title_bar(s, "三、通用规则提示")
box(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.2), fill=RGBColor(0x1A,0x28,0x42), line=GOLD, line_w=Pt(1), radius=True)
box(s, Inches(6.85), Inches(1.4), Inches(6.0), Inches(5.2), fill=RGBColor(0x1A,0x28,0x42), line=GOLD, line_w=Pt(1), radius=True)
txt(s, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.4), [("◆ 退货规则", 15, LGOLD, True)])
txt(s, Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4), [("◆ 收益核心", 15, LGOLD, True)])
rl=[("- 消费者：未使用、未损坏 1 个月内无条件退货",13,LIGHT,False,5),
    ("- 经销商：首次进货 1 月内可退；复购后仅质量问题可退",13,LIGHT,False,5),
    ("- 往返运费自理；违规乱价不予退货",13,LIGHT,False,5),
    ("- 退货后 1 年内不可重新参与销售推广",13,LIGHT,False,5)]
il=[("- 所有收益最终取决于自身团队架构与销售业绩",13,LIGHT,False,5),
    ("- 多劳多得，需踩准节点、跟上晋升节奏",13,LIGHT,False,5),
    ("- 月度计算周期：当月 26 日 — 次月 25 日",13,LIGHT,False,5)]
txt(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.4), rl)
txt(s, Inches(7.05), Inches(2.0), Inches(5.6), Inches(4.4), il)

# ---------- Slide 7: 封底 ----------
s = prs.slides.add_slide(BLANK); bg(s, DEEP)
box(s, 0, Inches(3.0), SW, Pt(3), fill=GOLD)
txt(s, Inches(1), Inches(3.2), SW-Inches(2), Inches(1.2),
    [("多劳多得 · 踩准节点 · 跟上节奏", 32, WHITE, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.5), SW-Inches(2), Inches(0.6),
    [("一四油经销商制度参考 · 内部培训资料", 14, LGOLD, False)], align=PP_ALIGN.CENTER)

prs.save("/workspace/经销商进阶计划.pptx")
print("PPT saved: /workspace/经销商进阶计划.pptx, slides =", len(prs.slides._sldIdLst))
